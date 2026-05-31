"""
PMtool Team Onboarding Presentation Generator
Run: python3 scripts/generate-onboarding-pptx.py
Output: /Users/chetanpatil/Desktop/PMtool_Onboarding.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1C, 0x22, 0x35)   # sidebar bg
INDIGO      = RGBColor(0x43, 0x61, 0xEE)   # primary
INDIGO_DARK = RGBColor(0x34, 0x51, 0xD1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xEE, 0xF0, 0xF7)
MID_GREY    = RGBColor(0x9D, 0xA8, 0xC3)
TEXT_DARK   = RGBColor(0x1A, 0x1F, 0x36)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)
VIOLET      = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # fully blank — we draw everything ourselves


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=TEXT_DARK, line_spacing_pt=4,
                  bold_first=False):
    """lines: list of str. Prefix with '•' for bullet."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txBox


def slide_header(slide, title, subtitle=None, accent=INDIGO):
    """Left navy sidebar + title block at top."""
    # Full background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_GREY)
    # Top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, fill_rgb=accent)
    # Left sidebar strip
    add_rect(slide, 0, 0, 0.06, 7.5, fill_rgb=accent)
    # Title
    add_text(slide, title, 0.4, 0.2, 12.5, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12.5, 0.4,
                 font_size=14, color=MID_GREY)


def role_badge(slide, label, left, top, bg_rgb, text_rgb=WHITE):
    add_rect(slide, left, top, 1.6, 0.35, fill_rgb=bg_rgb)
    add_text(slide, label, left + 0.05, top + 0.04, 1.5, 0.28,
             font_size=11, bold=True, color=text_rgb, align=PP_ALIGN.CENTER)


def card(slide, left, top, width, height, title, lines, accent_rgb=INDIGO,
         font_size=12, title_size=13):
    # Card background
    add_rect(slide, left, top, width, height, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE4, 0xE7, 0xF0), line_width_pt=0.75)
    # Left accent
    add_rect(slide, left, top, 0.05, height, fill_rgb=accent_rgb)
    # Title
    add_text(slide, title, left + 0.12, top + 0.1, width - 0.2, 0.3,
             font_size=title_size, bold=True, color=accent_rgb)
    # Body lines
    body_top = top + 0.42
    for line in lines:
        add_text(slide, line, left + 0.12, body_top, width - 0.2, 0.28,
                 font_size=font_size, color=TEXT_DARK)
        body_top += 0.28


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# Background
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
# Gradient hint — indigo rectangle top-right
add_rect(s, 7, 0, 6.33, 3.5, fill_rgb=RGBColor(0x2A, 0x34, 0x54))
# Bottom bar
add_rect(s, 0, 6.9, 13.33, 0.6, fill_rgb=INDIGO)

# Logo text
add_text(s, "PMtool", 0.6, 1.4, 6, 1.2,
         font_size=64, bold=True, color=WHITE)
add_text(s, "Production Management Platform", 0.6, 2.6, 8, 0.6,
         font_size=22, color=MID_GREY)

add_text(s, "Team Onboarding Guide", 0.6, 3.4, 8, 0.5,
         font_size=16, color=INDIGO, bold=True)

# Stats row
for i, (val, label) in enumerate([("13", "Workstations"), ("5", "User Roles"), ("Real-time", "Tracking")]):
    x = 0.6 + i * 2.6
    add_rect(s, x, 4.5, 2.2, 1.1, fill_rgb=RGBColor(0x25, 0x2D, 0x45))
    add_text(s, val, x + 0.1, 4.55, 2.0, 0.5, font_size=28, bold=True, color=INDIGO)
    add_text(s, label, x + 0.1, 5.0, 2.0, 0.35, font_size=12, color=MID_GREY)

add_text(s, "Confidential — Internal Use Only", 0.6, 7.1, 10, 0.28,
         font_size=10, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is PMtool?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "What is PMtool?", "A real-time production floor tracking system for energy meter manufacturing")

# Three value cards
for i, (icon, title, body) in enumerate([
    ("📋", "Track Every Meter",
     ["Each meter moves through 13", "quality checkpoints — tracked", "in real time from arrival", "to dispatch."]),
    ("👥", "Role-Based Access",
     ["5 roles: Super Admin, Admin,", "Supervisor, Operator, QA.", "Everyone sees exactly what", "they need — nothing more."]),
    ("📊", "Live Reports",
     ["Production throughput,", "failure history, tamper test", "results — filtered by date,", "station, or serial number."]),
]):
    x = 0.4 + i * 4.3
    add_rect(s, x, 1.5, 3.9, 4.6, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE4, 0xE7, 0xF0), line_width_pt=0.75)
    add_rect(s, x, 1.5, 3.9, 0.06, fill_rgb=INDIGO)
    add_text(s, icon, x + 0.15, 1.65, 0.6, 0.5, font_size=24)
    add_text(s, title, x + 0.15, 2.2, 3.6, 0.4, font_size=16, bold=True, color=INDIGO)
    for j, line in enumerate(body):
        add_text(s, line, x + 0.15, 2.7 + j * 0.35, 3.6, 0.35, font_size=13, color=TEXT_DARK)

# Bottom tagline
add_text(s, "No paper. No spreadsheets. No guessing.",
         0.4, 6.4, 12.5, 0.5, font_size=15, bold=True,
         color=INDIGO, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The 13-Stage Manufacturing Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "The 13-Stage Manufacturing Pipeline",
             "Every energy meter passes through all 13 quality checkpoints in sequence")

stages = [
    ("WS1",  "Incoming Inspection",          INDIGO),
    ("WS2",  "SMD / PCB Soldering",          INDIGO),
    ("WS3",  "PCBA Incoming",                INDIGO),
    ("WS4",  "Base Assembly",                INDIGO),
    ("WS5",  "Functional Testing",           VIOLET),
    ("WS6",  "Cover Assembly",               INDIGO),
    ("WS7",  "Error Compensation",           VIOLET),
    ("WS8",  "Tamper Test",                  AMBER),
    ("WS9",  "HV-IR Test",                   RED),
    ("WS10", "Soaking Test",                 INDIGO),
    ("WS11", "Final Testing",                VIOLET),
    ("WS12", "Sealing",                      INDIGO),
    ("WS13", "Packing & Dispatch",           GREEN),
]

cols = 7
rows = 2
for idx, (ws, name, color) in enumerate(stages):
    col = idx % cols
    row = idx // cols
    x = 0.35 + col * 1.85
    y = 1.55 + row * 2.2
    add_rect(s, x, y, 1.65, 1.8, fill_rgb=WHITE,
             line_rgb=color, line_width_pt=1.5)
    add_rect(s, x, y, 1.65, 0.32, fill_rgb=color)
    add_text(s, ws, x + 0.05, y + 0.04, 1.55, 0.26,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Word-wrap name into 2 lines
    words = name.split()
    if len(words) <= 2:
        lines = [name]
    else:
        mid = len(words) // 2
        lines = [" ".join(words[:mid]), " ".join(words[mid:])]
    for li, line in enumerate(lines):
        add_text(s, line, x + 0.05, y + 0.42 + li * 0.35, 1.55, 0.35,
                 font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER, bold=True)

# Legend
for color, label in [(VIOLET, "Testing / QC"), (AMBER, "Tamper"), (RED, "HV-IR"), (GREEN, "Dispatch")]:
    pass  # skipped for brevity — colours are self-explanatory

add_text(s, "🟣 Testing/QC   🟡 Tamper   🔴 HV-IR   🟢 Final Dispatch",
         0.35, 6.55, 12.5, 0.4, font_size=11, color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Roles & Permissions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Roles & Permissions", "Five roles — each with a distinct scope of access")

role_data = [
    ("Admin",      INDIGO, "Day-to-day management",
     ["Create & manage users", "Reset operator passwords", "View all reports & audit log", "Assign operators to stations", "Full team visibility"]),
    ("Supervisor", VIOLET, "Floor oversight",
     ["Monitor all 13 workstations", "Reset operator passwords", "View reports & audit log", "Cannot create new users", "Real-time queue visibility"]),
    ("Operator",   GREEN,  "Production work",
     ["See assigned workstation queue", "Submit quality check results", "Tag meters for rework", "Add comments to submissions", "No admin access"]),
    ("QA",         AMBER,  "Quality review",
     ["Read-only access to all data", "View stage history per meter", "Access all 3 report types", "Cannot modify any data", "Audit trail visibility"]),
]

for i, (role, color, tagline, bullets) in enumerate(role_data):
    x = 0.4 + i * 3.15
    y = 1.4
    add_rect(s, x, y, 2.9, 5.4, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE4, 0xE7, 0xF0), line_width_pt=0.75)
    add_rect(s, x, y, 2.9, 0.5, fill_rgb=color)
    add_text(s, role, x + 0.08, y + 0.1, 2.74, 0.35,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, tagline, x + 0.08, y + 0.62, 2.74, 0.38,
             font_size=10, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        add_text(s, f"✓  {b}", x + 0.12, y + 1.12 + j * 0.78, 2.66, 0.68,
                 font_size=11, color=TEXT_DARK)

add_text(s, "Login URL: /login (Operators)  ·  /admin/login (Admin & Supervisor)",
         0.4, 6.88, 12.5, 0.35, font_size=11, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — How to Log In
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "How to Log In", "Two login pages depending on your role")

# Operator path
add_rect(s, 0.4, 1.4, 5.8, 4.8, fill_rgb=WHITE,
         line_rgb=GREEN, line_width_pt=1.5)
add_rect(s, 0.4, 1.4, 5.8, 0.45, fill_rgb=GREEN)
add_text(s, "Operators", 0.5, 1.46, 5.6, 0.35,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

op_steps = [
    ("1", "Open your browser"),
    ("2", "Go to  pmtool-3f8db.web.app/login"),
    ("3", "Enter your email & password"),
    ("4", "You'll land on your assigned workstation"),
    ("5", "Start processing your queue!"),
]
for j, (num, step) in enumerate(op_steps):
    y = 2.05 + j * 0.72
    add_rect(s, 0.65, y, 0.38, 0.38, fill_rgb=GREEN)
    add_text(s, num, 0.65, y + 0.02, 0.38, 0.34,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, 1.15, y + 0.04, 4.7, 0.34, font_size=13, color=TEXT_DARK)

# Admin path
add_rect(s, 7.0, 1.4, 5.8, 4.8, fill_rgb=WHITE,
         line_rgb=INDIGO, line_width_pt=1.5)
add_rect(s, 7.0, 1.4, 5.8, 0.45, fill_rgb=INDIGO)
add_text(s, "Admin & Supervisor", 7.1, 1.46, 5.6, 0.35,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

admin_steps = [
    ("1", "Open your browser"),
    ("2", "Go to  pmtool-3f8db.web.app/admin/login"),
    ("3", "Enter your email & password"),
    ("4", "You'll land on the Dashboard"),
    ("5", "Forgot password? Use the link on that page"),
]
for j, (num, step) in enumerate(admin_steps):
    y = 2.05 + j * 0.72
    add_rect(s, 7.25, y, 0.38, 0.38, fill_rgb=INDIGO)
    add_text(s, num, 7.25, y + 0.02, 0.38, 0.34,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, 7.75, y + 0.04, 4.7, 0.34, font_size=13, color=TEXT_DARK)

add_text(s, "⚠  Always open in an Incognito window if you were previously logged in as someone else",
         0.4, 6.35, 12.5, 0.45, font_size=11, color=AMBER, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Operator Workflow
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Operator Workflow", "What you do at your workstation — step by step")

steps = [
    (GREEN,  "STEP 1", "Log in",
     ["Go to /login", "Enter your credentials", "You land on your station's queue"]),
    (INDIGO, "STEP 2", "Pick a meter",
     ["You see a list of meters waiting", "Click on any meter to start work", "Serial number is shown"]),
    (VIOLET, "STEP 3", "Fill the checklist",
     ["Each parameter has Pass / Fail", "Add a reading where required", "All fields must be filled"]),
    (AMBER,  "STEP 4", "Submit",
     ["All Pass → meter moves to next WS", "Any Fail → choose Rework stage", "Add a comment explaining why"]),
    (RED,    "STEP 5", "Rework handling",
     ["Rework meter goes back to target WS", "Counter tracks how many reworks", "QA can review full history"]),
]

for i, (color, label, title, bullets) in enumerate(steps):
    x = 0.3 + i * 2.58
    add_rect(s, x, 1.35, 2.38, 5.4, fill_rgb=WHITE,
             line_rgb=color, line_width_pt=1.5)
    add_rect(s, x, 1.35, 2.38, 0.32, fill_rgb=color)
    add_text(s, label, x + 0.05, 1.37, 2.28, 0.28,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.1, 1.77, 2.18, 0.4,
             font_size=14, bold=True, color=color)
    for j, b in enumerate(bullets):
        add_text(s, f"→  {b}", x + 0.1, 2.3 + j * 0.75, 2.18, 0.65,
                 font_size=11, color=TEXT_DARK)

add_text(s, "The queue is always live — no refresh needed. New meters appear automatically.",
         0.3, 6.85, 12.5, 0.35, font_size=11, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Admin & Supervisor Workflow
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Admin & Supervisor Tasks", "Managing the team and monitoring the floor")

tasks = [
    ("👤  Add a new user",
     ["Go to Team → click Add User", "Enter name, email, role, temp password", "User can log in immediately"]),
    ("🔑  Reset a password",
     ["Team page → Reset Password button", "System generates a new temp password", "Share it securely with the operator"]),
    ("🏭  Assign operator to workstation",
     ["Team page → Edit user", "Select their workstation", "They will see that station's queue"]),
    ("📊  Run a report",
     ["Reports → choose tab (Production /", "Failure / Tamper)", "Set date range → export CSV"]),
    ("📋  Audit log",
     ["Every delete & password reset", "is recorded automatically", "Filter by user, date, or action"]),
    ("🖥️  Workstation monitor",
     ["See all 13 stations at a glance", "Machine status: Green / Yellow / Red", "Queue depth per station"]),
]

cols = 3
for i, (title, bullets) in enumerate(tasks):
    col = i % cols
    row = i // cols
    x = 0.3 + col * 4.35
    y = 1.4 + row * 2.7
    add_rect(s, x, y, 4.1, 2.45, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE4, 0xE7, 0xF0), line_width_pt=0.75)
    add_rect(s, x, y, 0.06, 2.45, fill_rgb=INDIGO)
    add_text(s, title, x + 0.15, y + 0.1, 3.85, 0.4,
             font_size=13, bold=True, color=INDIGO)
    for j, b in enumerate(bullets):
        add_text(s, f"•  {b}", x + 0.15, y + 0.58 + j * 0.55, 3.85, 0.5,
                 font_size=11, color=TEXT_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Reports Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Reports", "Three built-in reports — all filterable and exportable to CSV")

reports = [
    (INDIGO, "📈  Production Report",
     "What it shows",
     ["Total meters completed per day", "Breakdown by meter type", "Throughput by workstation"],
     "Use it for",
     ["Daily production targets", "Capacity planning", "Client delivery commitments"]),
    (RED,    "⚠️  Failure History",
     "What it shows",
     ["All rework events with reason", "Failed parameters per meter", "Which stages fail most often"],
     "Use it for",
     ["Root cause analysis", "Training gaps", "Process improvement"]),
    (AMBER,  "🔒  Tamper Test",
     "What it shows",
     ["Pass / Fail per tamper parameter", "Magnetic, Neutral, Cover, Earth", "Result per serial number"],
     "Use it for",
     ["Compliance records", "Dispute resolution", "Audit evidence"]),
]

for i, (color, title, h1, l1, h2, l2) in enumerate(reports):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.35, 4.1, 5.5, fill_rgb=WHITE,
             line_rgb=color, line_width_pt=1.5)
    add_rect(s, x, 1.35, 4.1, 0.5, fill_rgb=color)
    add_text(s, title, x + 0.1, 1.38, 3.9, 0.42,
             font_size=14, bold=True, color=WHITE)

    add_text(s, h1, x + 0.12, 1.95, 3.86, 0.3,
             font_size=11, bold=True, color=color)
    for j, b in enumerate(l1):
        add_text(s, f"→  {b}", x + 0.12, 2.3 + j * 0.45, 3.86, 0.42, font_size=11, color=TEXT_DARK)

    add_text(s, h2, x + 0.12, 3.75, 3.86, 0.3,
             font_size=11, bold=True, color=color)
    for j, b in enumerate(l2):
        add_text(s, f"→  {b}", x + 0.12, 4.1 + j * 0.45, 3.86, 0.42, font_size=11, color=TEXT_DARK)

add_text(s, "Tip: Use the Quick Select presets (Today / 7 days / 30 days) to jump to common date ranges instantly.",
         0.3, 6.9, 12.5, 0.35, font_size=11, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Do's & Don'ts
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Do's & Don'ts", "Guidelines for using PMtool correctly")

dos = [
    "Log out when you leave your station",
    "Always use your own login — no sharing",
    "Mark Rework with an accurate comment",
    "Contact your admin if you forget your password",
    "Use the correct serial number from the physical meter",
    "Report any system errors to your supervisor immediately",
]
donts = [
    "Don't share your password with anyone",
    "Don't mark a meter as Pass if any check failed",
    "Don't use another operator's login to 'save time'",
    "Don't close the browser mid-submission",
    "Don't guess serial numbers — scan or copy carefully",
    "Don't delete data — contact admin if data looks wrong",
]

# Do's card
add_rect(s, 0.4, 1.35, 5.9, 5.5, fill_rgb=RGBColor(0xF0, 0xFD, 0xF4),
         line_rgb=GREEN, line_width_pt=1.5)
add_rect(s, 0.4, 1.35, 5.9, 0.45, fill_rgb=GREEN)
add_text(s, "✓  Do's", 0.55, 1.4, 5.6, 0.35,
         font_size=16, bold=True, color=WHITE)
for j, item in enumerate(dos):
    add_text(s, f"✓  {item}", 0.55, 1.95 + j * 0.72, 5.6, 0.62,
             font_size=12, color=RGBColor(0x16, 0x6D, 0x3B))

# Don'ts card
add_rect(s, 7.0, 1.35, 5.9, 5.5, fill_rgb=RGBColor(0xFF, 0xF1, 0xF2),
         line_rgb=RED, line_width_pt=1.5)
add_rect(s, 7.0, 1.35, 5.9, 0.45, fill_rgb=RED)
add_text(s, "✗  Don'ts", 7.15, 1.4, 5.6, 0.35,
         font_size=16, bold=True, color=WHITE)
for j, item in enumerate(donts):
    add_text(s, f"✗  {item}", 7.15, 1.95 + j * 0.72, 5.6, 0.62,
             font_size=12, color=RGBColor(0x9B, 0x1C, 0x1C))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FAQ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Frequently Asked Questions")

faqs = [
    ("I forgot my password — what do I do?",
     "If you're an operator, contact your admin or supervisor. They can reset it from the Team page and give you a new temporary password. If you're an admin, use the 'Forgot password?' link on /admin/login."),
    ("The meter I'm looking for isn't in my queue — where is it?",
     "It may be at a different workstation, still in rework at a previous stage, or already completed. Ask your supervisor to check the meter's serial number in the system."),
    ("I submitted wrong results — can I undo?",
     "You cannot undo a submission. Immediately inform your supervisor. They can review the audit trail and the QA team can note the discrepancy."),
    ("The system is slow or not loading — what should I do?",
     "Refresh the page (Ctrl+R / Cmd+R). If the issue persists, report it to your admin with a screenshot of any error messages."),
    ("Can I use PMtool on my phone?",
     "Yes, the interface is responsive and works on mobile browsers. For best experience use Chrome on Android or Safari on iPhone."),
    ("Why does my screen show a blank dashboard after login?",
     "This usually happens when a stale session is cached. Open a new incognito / private browser window and log in fresh."),
]

for i, (q, a) in enumerate(faqs):
    row = i % 3
    col = i // 3
    x = 0.35 + col * 6.65
    y = 1.35 + row * 1.9
    add_rect(s, x, y, 6.3, 1.72, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE4, 0xE7, 0xF0), line_width_pt=0.75)
    add_rect(s, x, y, 0.05, 1.72, fill_rgb=INDIGO)
    add_text(s, f"Q: {q}", x + 0.12, y + 0.08, 6.05, 0.4,
             font_size=11, bold=True, color=INDIGO)
    add_text(s, a, x + 0.12, y + 0.52, 6.05, 1.1,
             font_size=10, color=TEXT_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Contact & Support
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
add_rect(s, 0, 0, 13.33, 0.06, fill_rgb=INDIGO)
add_rect(s, 0, 6.8, 13.33, 0.7, fill_rgb=INDIGO)

add_text(s, "Need Help?", 0.6, 0.8, 12, 0.8,
         font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Your first point of contact is always your Supervisor or Admin.",
         0.6, 1.6, 12, 0.5, font_size=16, color=MID_GREY, align=PP_ALIGN.CENTER)

contacts = [
    ("🔑", "Password issues",    "Contact your Admin → Team page → Reset Password"),
    ("🖥️",  "Station / queue",   "Contact your Supervisor on the floor"),
    ("📊", "Report access",      "Contact your Admin to check your role permissions"),
    ("🔧", "System / technical", "Contact the PMtool system administrator"),
]

for i, (icon, title, detail) in enumerate(contacts):
    x = 0.5 + i * 3.1
    add_rect(s, x, 2.4, 2.8, 2.8, fill_rgb=RGBColor(0x25, 0x2D, 0x45))
    add_text(s, icon,   x + 0.1, 2.5,  2.6, 0.55, font_size=26, align=PP_ALIGN.CENTER)
    add_text(s, title,  x + 0.1, 3.15, 2.6, 0.4,  font_size=13, bold=True,
             color=INDIGO, align=PP_ALIGN.CENTER)
    add_text(s, detail, x + 0.1, 3.65, 2.6, 1.3,  font_size=10,
             color=MID_GREY, align=PP_ALIGN.CENTER)

add_text(s, "PMtool  ·  Production Management Platform  ·  Confidential",
         0.6, 7.08, 12, 0.3, font_size=10, color=WHITE,
         align=PP_ALIGN.CENTER, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────────
OUTPUT = "/Users/chetanpatil/Desktop/PMtool_Onboarding.pptx"
prs.save(OUTPUT)
print(f"\n✅  Presentation saved to:\n   {OUTPUT}\n")
print("  11 slides:")
print("   1  Cover")
print("   2  What is PMtool?")
print("   3  The 13-Stage Pipeline")
print("   4  Roles & Permissions")
print("   5  How to Log In")
print("   6  Operator Workflow")
print("   7  Admin & Supervisor Tasks")
print("   8  Reports Overview")
print("   9  Do's & Don'ts")
print("  10  FAQ")
print("  11  Contact & Support")
