#!/usr/bin/env python3
"""
PMtool Executive/Stakeholder PowerPoint Deck Generator
Generates a 10-slide presentation saved to ~/Desktop/PMtool_Executive_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── Brand colors ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1C, 0x22, 0x35)
INDIGO  = RGBColor(0x43, 0x61, 0xEE)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
LGREY   = RGBColor(0xEE, 0xF0, 0xF7)
MGREY   = RGBColor(0x9D, 0xA8, 0xC3)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# Tinted variants (card backgrounds)
RED_TINT    = RGBColor(0xFE, 0xE2, 0xE2)
INDIGO_TINT = RGBColor(0xE0, 0xE7, 0xFF)
VIOLET_TINT = RGBColor(0xED, 0xE9, 0xFE)
GREEN_TINT  = RGBColor(0xDC, 0xFC, 0xE7)
AMBER_TINT  = RGBColor(0xFE, 0xF3, 0xC7)

OUTPUT_PATH = os.path.expanduser("~/Desktop/PMtool_Executive_Deck.pptx")

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb and line_pt > 0:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, font_size,
             bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = TEXT_DARK
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def blank_slide(prs):
    """Add a slide using the blank layout (index 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def full_bg(slide, color):
    add_rect(slide, 0, 0, 13.33, 7.5, color)


def light_bg_with_sidebar(slide):
    """Standard light-grey background + indigo left sidebar."""
    full_bg(slide, LGREY)
    add_rect(slide, 0, 0, 0.25, 7.5, INDIGO)


def slide_title(slide, title_text, top=0.35, left=0.55, width=12.0):
    add_text(slide, title_text, left, top, width, 0.7,
             font_size=28, bold=True, color=TEXT_DARK)


# ─── Slide builders ──────────────────────────────────────────────────────────

def build_slide1(prs):
    """Cover – dark navy."""
    s = blank_slide(prs)
    full_bg(s, NAVY)

    # Top accent bar
    add_rect(s, 0, 0, 13.33, 0.18, INDIGO)

    # PMtool wordmark
    add_text(s, "PMtool", 1.2, 1.8, 8.0, 1.4,
             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Sub-title
    add_text(s, "Production Management Platform", 1.2, 3.1, 9.0, 0.55,
             font_size=20, color=MGREY, align=PP_ALIGN.LEFT)

    # Tagline
    add_text(s,
             "Built for energy meter manufacturers who demand zero-defect delivery",
             1.2, 3.8, 10.5, 0.55,
             font_size=14, color=INDIGO, italic=True, align=PP_ALIGN.LEFT)

    # Bottom bar
    add_rect(s, 0, 6.85, 13.33, 0.65, INDIGO)
    add_text(s, "Executive Overview — 2026", 0.35, 6.9, 12.0, 0.5,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide2(prs):
    """The Problem We Solved."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "The Problem We Solved")

    # Horizontal rule under title
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    problems = [
        ("No Traceability",
         "Paper checklists, lost records, no way to trace a defect back to a station or operator"),
        ("Rework Black Hole",
         "When a meter fails, no one knows how many times it has been sent back or why"),
        ("Zero Visibility",
         "Managers had no real-time view of production status, bottlenecks, or throughput"),
    ]

    card_w = 3.7
    card_h = 3.8
    gap = 0.32
    start_x = 0.55
    top = 1.2

    for i, (title, body) in enumerate(problems):
        x = start_x + i * (card_w + gap)
        # card bg
        add_rect(s, x, top, card_w, card_h, RED_TINT, RED, 1.2)
        # red accent top bar
        add_rect(s, x, top, card_w, 0.14, RED)
        # number circle background
        add_rect(s, x + 0.15, top + 0.28, 0.42, 0.42, RED)
        add_text(s, str(i + 1), x + 0.15, top + 0.28, 0.42, 0.42,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        add_text(s, title, x + 0.1, top + 0.82, card_w - 0.2, 0.5,
                 font_size=15, bold=True, color=TEXT_DARK)
        # body
        add_text(s, body, x + 0.1, top + 1.4, card_w - 0.2, 2.2,
                 font_size=11, color=TEXT_DARK)


def build_slide3(prs):
    """The Solution."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "The Solution")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    # Central statement
    add_rect(s, 0.55, 1.1, 11.5, 0.85, INDIGO)
    add_text(s,
             "One platform. 13 checkpoints. Every meter. Every operator. Real time.",
             0.65, 1.12, 11.2, 0.8,
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    solutions = [
        ("Digital Quality Gates",
         "Every meter scanned through 13 mandatory QC stages. No skipping. No paper."),
        ("Live Floor Dashboard",
         "Managers see every workstation status, queue depth, and machine health right now."),
        ("Instant Audit Trail",
         "Every submission, failure, rework, and password reset is logged permanently."),
    ]

    card_w = 3.7
    card_h = 3.7
    gap = 0.32
    start_x = 0.55
    top = 2.15

    for i, (title, body) in enumerate(solutions):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, card_h, INDIGO_TINT, INDIGO, 1.2)
        add_rect(s, x, top, card_w, 0.14, INDIGO)
        add_rect(s, x + 0.15, top + 0.28, 0.42, 0.42, INDIGO)
        add_text(s, str(i + 1), x + 0.15, top + 0.28, 0.42, 0.42,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, x + 0.1, top + 0.82, card_w - 0.2, 0.5,
                 font_size=15, bold=True, color=INDIGO)
        add_text(s, body, x + 0.1, top + 1.4, card_w - 0.2, 2.1,
                 font_size=11, color=TEXT_DARK)


def build_slide4(prs):
    """13-Stage Quality Pipeline."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)

    slide_title(s, "13-Stage Quality Pipeline")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)
    add_text(s,
             "Every energy meter travels this path — digitally tracked at each handover",
             0.55, 1.02, 11.5, 0.38,
             font_size=12, color=MGREY, italic=True)

    stages = [
        ("WS1",  "Incoming\nInspection",  INDIGO),
        ("WS2",  "SMD\nSoldering",        VIOLET),
        ("WS3",  "PCBA\nIncoming",        INDIGO),
        ("WS4",  "Base\nAssembly",        VIOLET),
        ("WS5",  "Functional\nTesting",   AMBER),
        ("WS6",  "Cover\nAssembly",       INDIGO),
        ("WS7",  "Error\nCompensation",   VIOLET),
        ("WS8",  "Tamper\nTest",          RED),
        ("WS9",  "HV-IR\nTest",           RED),
        ("WS10", "Soaking\nTest",         AMBER),
        ("WS11", "Final\nTesting",        GREEN),
        ("WS12", "Sealing",              INDIGO),
        ("WS13", "Packing &\nDispatch",   GREEN),
    ]

    tile_w = 1.68
    tile_h = 1.2
    gap_x = 0.1
    gap_y = 0.18
    row1_count = 7
    start_x = 0.38
    row1_y = 1.55
    row2_y = row1_y + tile_h + gap_y

    for i, (ws, label, color) in enumerate(stages):
        if i < row1_count:
            x = start_x + i * (tile_w + gap_x)
            y = row1_y
        else:
            j = i - row1_count
            # center 6 tiles in row 2
            total_row2 = len(stages) - row1_count
            row2_total_w = total_row2 * tile_w + (total_row2 - 1) * gap_x
            row2_start_x = (13.33 - row2_total_w) / 2
            x = row2_start_x + j * (tile_w + gap_x)
            y = row2_y

        add_rect(s, x, y, tile_w, tile_h, color)
        # WS label
        add_text(s, ws, x, y + 0.06, tile_w, 0.3,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # stage name
        add_text(s, label, x, y + 0.36, tile_w, 0.75,
                 font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide5(prs):
    """Key Capabilities — 3×2 grid."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "Key Capabilities")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    caps = [
        ("Role-Based Access",
         "4 roles (Admin, Supervisor, Operator, QA). Each person sees only what they need.",
         INDIGO, INDIGO_TINT),
        ("Real-Time Queue",
         "Operators see their live work queue. No refresh. No delay.",
         VIOLET, VIOLET_TINT),
        ("Rework Tracking",
         "Every failure is logged with reason, stage, and operator. Rework counter per meter.",
         RED, RED_TINT),
        ("Reports & Export",
         "Production throughput, failure history, tamper test results. Exportable to CSV.",
         GREEN, GREEN_TINT),
        ("Tamper Test Compliance",
         "Dedicated tamper test stage with per-parameter pass/fail and permanent record.",
         AMBER, AMBER_TINT),
        ("Immutable Audit Log",
         "Delete, password reset, every admin action — permanently recorded. Cannot be altered.",
         INDIGO, INDIGO_TINT),
    ]

    cols = 3
    card_w = 3.7
    card_h = 2.45
    gap_x = 0.27
    gap_y = 0.22
    start_x = 0.48
    start_y = 1.15

    for i, (title, body, accent, bg) in enumerate(caps):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rect(s, x, y, card_w, card_h, bg, accent, 1.2)
        add_rect(s, x, y, card_w, 0.12, accent)
        add_text(s, title, x + 0.12, y + 0.22, card_w - 0.24, 0.4,
                 font_size=13, bold=True, color=accent)
        add_text(s, body, x + 0.12, y + 0.68, card_w - 0.24, 1.6,
                 font_size=10, color=TEXT_DARK)


def build_slide6(prs):
    """Business Impact."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "Business Impact")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    metrics = [
        ("13",        "Quality checkpoints,\nzero paper",            INDIGO, INDIGO_TINT),
        ("100%",      "Traceability from\narrival to dispatch",       GREEN,  GREEN_TINT),
        ("Real-time", "Production visibility\nfor managers",          VIOLET, VIOLET_TINT),
        ("Permanent", "Audit trail, tamper\n& compliance records",    AMBER,  AMBER_TINT),
    ]

    tile_w = 5.45
    tile_h = 2.4
    gap_x = 0.35
    gap_y = 0.25
    start_x = 0.6
    start_y = 1.15

    for i, (big, label, accent, bg) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = start_x + col * (tile_w + gap_x)
        y = start_y + row * (tile_h + gap_y)
        add_rect(s, x, y, tile_w, tile_h, bg, accent, 1.5)
        add_rect(s, x, y, 0.14, tile_h, accent)
        add_text(s, big, x + 0.28, y + 0.18, tile_w - 0.38, 1.0,
                 font_size=38, bold=True, color=accent, align=PP_ALIGN.LEFT)
        add_text(s, label, x + 0.28, y + 1.18, tile_w - 0.38, 1.0,
                 font_size=12, color=TEXT_DARK, align=PP_ALIGN.LEFT)

    # Tagline below tiles
    add_text(s,
             "From clipboard to cloud — the same rigor, with proof.",
             0.55, 6.5, 11.5, 0.45,
             font_size=14, bold=True, color=INDIGO, italic=True, align=PP_ALIGN.CENTER)


def build_slide7(prs):
    """Who Uses PMtool."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "Who Uses PMtool")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    roles = [
        ("Admin",      INDIGO, INDIGO_TINT,
         "Manages team, resets passwords, views all reports"),
        ("Supervisor", VIOLET, VIOLET_TINT,
         "Monitors floor in real time, oversees rework"),
        ("Operator",   GREEN,  GREEN_TINT,
         "Processes meter queue, submits QC results"),
        ("QA",         AMBER,  AMBER_TINT,
         "Reviews history, accesses all reports, read-only"),
    ]

    card_w = 2.8
    card_h = 4.5
    gap = 0.3
    start_x = 0.55
    top = 1.25

    for i, (role, accent, bg, desc) in enumerate(roles):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, card_h, bg, accent, 1.5)
        add_rect(s, x, top, card_w, 0.28, accent)
        # Role icon circle
        add_rect(s, x + card_w / 2 - 0.42, top + 0.5, 0.84, 0.84, accent)
        add_text(s, role[0], x + card_w / 2 - 0.42, top + 0.5, 0.84, 0.84,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, role, x + 0.1, top + 1.55, card_w - 0.2, 0.5,
                 font_size=16, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + 0.15, top + 2.15, card_w - 0.3, 2.0,
                 font_size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def build_slide8(prs):
    """Technical Reliability."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "Built to Be Trusted")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)
    add_text(s, "Enterprise-grade infrastructure, zero maintenance burden",
             0.55, 1.0, 11.5, 0.38,
             font_size=12, color=MGREY, italic=True)

    points = [
        ("Cloud-Native",
         "Hosted on Google Firebase + Cloud Run. No servers. No maintenance. Auto-scaling."),
        ("Always Available",
         "99.95% Firebase SLA. Data replicated across Google's global infrastructure."),
        ("Secure by Design",
         "Role enforcement at the token level (not just the UI). Cookies HttpOnly. "
         "Audit log write-once at the database layer."),
        ("Real-Time Sync",
         "Firestore live listeners mean the floor updates instantly — "
         "no polling, no refresh, no stale data."),
    ]

    row_h = 1.18
    gap = 0.16
    start_y = 1.5
    for i, (title, body) in enumerate(points):
        y = start_y + i * (row_h + gap)
        add_rect(s, 0.55, y, 11.5, row_h, WHITE, MGREY, 0.6)
        add_rect(s, 0.55, y, 0.22, row_h, INDIGO)
        add_text(s, title, 0.92, y + 0.1, 4.5, 0.42,
                 font_size=14, bold=True, color=INDIGO)
        add_text(s, body, 0.92, y + 0.55, 10.9, 0.6,
                 font_size=11, color=TEXT_DARK)


def build_slide9(prs):
    """Roadmap."""
    s = blank_slide(prs)
    light_bg_with_sidebar(s)
    slide_title(s, "What's Coming Next")
    add_rect(s, 0.55, 0.95, 11.5, 0.04, INDIGO)

    phases = [
        ("Phase 1",
         "DONE",
         GREEN, GREEN_TINT,
         "RBAC auth, 13-stage pipeline, rework tracking, reports, audit log, tamper test"),
        ("Phase 2",
         "NEXT",
         INDIGO, INDIGO_TINT,
         "Shift management, barcode/QR scan input, WhatsApp/SMS alerts for rework spikes"),
        ("Phase 3",
         "FUTURE",
         MGREY, LGREY,
         "Predictive defect analytics, supplier scorecard, mobile-native app"),
    ]

    card_w = 3.7
    card_h = 4.2
    gap = 0.32
    start_x = 0.55
    top = 1.3

    for i, (phase, status, accent, bg, items) in enumerate(phases):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, card_h, bg, accent, 1.5)
        add_rect(s, x, top, card_w, 0.55, accent)
        add_text(s, phase, x + 0.12, top + 0.04, card_w - 0.24, 0.32,
                 font_size=14, bold=True, color=WHITE)
        # Status badge
        badge_x = x + card_w - 1.05
        badge_y = top + 0.1
        add_rect(s, badge_x, badge_y, 0.9, 0.3, WHITE)
        add_text(s, status, badge_x, badge_y, 0.9, 0.3,
                 font_size=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(s, items, x + 0.15, top + 0.75, card_w - 0.3, 3.3,
                 font_size=11, color=TEXT_DARK)

    # Connector arrows between cards (simple rect lines)
    for i in range(2):
        arrow_x = start_x + (i + 1) * (card_w + gap) - gap / 2 - 0.06
        add_rect(s, arrow_x, top + card_h / 2 - 0.06, 0.32, 0.12, MGREY)


def build_slide10(prs):
    """Closing / Thank You."""
    s = blank_slide(prs)
    full_bg(s, NAVY)
    add_rect(s, 0, 0, 13.33, 0.18, INDIGO)

    add_text(s, "Thank You", 1.0, 1.8, 11.0, 1.6,
             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "PMtool — Production Management Platform",
             1.0, 3.55, 11.0, 0.7,
             font_size=20, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    add_text(s, "Questions? Contact your implementation team.",
             1.0, 4.4, 11.0, 0.55,
             font_size=14, color=MGREY, align=PP_ALIGN.CENTER)

    add_rect(s, 0, 6.85, 13.33, 0.65, INDIGO)
    add_text(s, "Confidential — For client use only",
             0.35, 6.9, 12.5, 0.5,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    build_slide8(prs)
    build_slide9(prs)
    build_slide10(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
